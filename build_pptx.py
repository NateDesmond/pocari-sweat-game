#!/usr/bin/env python3
"""Generate a branded Pocari Sweat history PowerPoint (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Brand palette ----
BLUE  = RGBColor(0x00, 0x9F, 0xE3)
DEEP  = RGBColor(0x00, 0x66, 0xB3)
NAVY  = RGBColor(0x00, 0x35, 0x5C)
NIGHT = RGBColor(0x04, 0x11, 0x1F)
PALE  = RGBColor(0xEA, 0xF8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE   = RGBColor(0xAE, 0xF0, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def grad(shape, c1, c2, angle=90):
    """Apply a two-stop linear gradient fill to a shape via raw XML."""
    sp = shape.fill._xPr  # spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        e = sp.find(qn(tag))
        if e is not None:
            sp.remove(e)
    g = sp.makeelement(qn('a:gradFill'), {})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = g.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = g.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(clr)
        lst.append(gs)
    g.append(lst)
    lin = g.makeelement(qn('a:lin'), {'ang': str(int(angle * 60000)), 'scaled': '1'})
    g.append(lin)
    # insert gradient before line element if present
    ln = sp.find(qn('a:ln'))
    sp.insert(list(sp).index(ln), g) if ln is not None else sp.append(g)


def bg(slide, kind):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    rect.line.fill.background()
    rect.shadow.inherit = False
    if kind == 'light':
        grad(rect, (0xEA, 0xF8, 0xFF), (0x8F, 0xD4, 0xF3), 145)
    elif kind == 'dark':
        grad(rect, (0x00, 0x3A, 0x66), (0x00, 0x9F, 0xE3), 150)
    elif kind == 'night':
        grad(rect, (0x10, 0x40, 0x6B), (0x04, 0x11, 0x1F), 150)
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def setrun(r, text, size, color, bold=True, italic=False, font='Helvetica', spacing=None):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(spacing * 100)))


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def eyebrow(tf, text, color):
    p = tf.paragraphs[0]
    setrun(p.add_run(), text.upper(), 13, color, bold=True, spacing=3)


def logo(slide, on_dark):
    tf = tb(slide, Inches(0.45), Inches(0.28), Inches(4), Inches(0.5))
    p = tf.paragraphs[0]
    setrun(p.add_run(), "POCARI", 16, WHITE if on_dark else DEEP, bold=True, italic=True)
    setrun(p.add_run(), "SWEAT", 16, ICE if on_dark else BLUE, bold=True, italic=True)


def bullets(tf, items, color, dark=False):
    for it in items:
        p = tf.add_paragraph()
        p.space_after = Pt(10)
        setrun(p.add_run(), ("●  " if dark else "💧  "), 18, BLUE if dark else DEEP, bold=True)
        setrun(p.add_run(), it, 18, color, bold=False)


# ---------- SLIDES ----------
def s_title():
    s = prs.slides.add_slide(BLANK); bg(s, 'dark'); logo(s, True)
    tf = tb(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(4.6))
    eyebrow(tf, "A History of an Icon", ICE)
    p = tf.add_paragraph(); p.space_before = Pt(14)
    setrun(p.add_run(), "POCARI SWEAT", 72, WHITE, bold=True, italic=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(20)
    setrun(p2.add_run(), "How a pharmaceutical company turned the science of an IV drip into "
           "the drink you can sweat — and a billion-dollar cultural icon.", 22, ICE, bold=False)
    p3 = tf.add_paragraph(); p3.space_before = Pt(22)
    setrun(p3.add_run(), "Otsuka Pharmaceutical  ·  Launched 1980  ·  Tokushima, Japan", 14, ICE, bold=True)


def content(kind, num, eb, head, lead_parts, items=None, items_dark=False, stats=None):
    s = prs.slides.add_slide(BLANK); bg(s, kind)
    on_dark = kind in ('dark', 'night')
    logo(s, on_dark)
    txt = WHITE if on_dark else NAVY
    accent = ICE if on_dark else DEEP
    tf = tb(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(5.6))
    eyebrow(tf, f"{num} · {eb}", accent)
    p = tf.add_paragraph(); p.space_before = Pt(10)
    setrun(p.add_run(), head, 40, txt, bold=True, italic=False)
    lp = tf.add_paragraph(); lp.space_before = Pt(18)
    for seg, strong in lead_parts:
        setrun(lp.add_run(), seg, 21, accent if strong else txt, bold=strong)
    if stats:
        sp = tf.add_paragraph(); sp.space_before = Pt(20)
        for j, (n, l) in enumerate(stats):
            setrun(sp.add_run(), ("        " if j else "") + n + "  ", 30, accent, bold=True, italic=True)
            setrun(sp.add_run(), l, 14, txt, bold=True)
    if items:
        bullets(tf, items, txt, dark=on_dark)
    return s


s_title()

content('light', '01', 'The Maker',
        "Born from a medicine company, not a soda company.",
        [("Pocari Sweat comes from ", False), ("Otsuka Pharmaceutical", True),
         (" — with roots tracing to ", False), ("1921", True),
         (", and the pharmaceutical company established in ", False), ("1964", True), (".", False)],
        items=["Decades of expertise making intravenous (IV) solutions for hospitals",
               "That fluid-and-electrolyte know-how became the base for a consumer drink"])

content('night', '02', 'The Spark',
        "A lightbulb moment in a Mexican hospital.",
        [("In the early 1970s, researcher ", False), ("Rokuro Harima", True),
         (" was hospitalized with diarrhea on a trip to Mexico. Watching a doctor drink an ", False),
         ("IV rehydration solution", True), (", he wondered: what if you could bottle that?", False)],
        items=["The idea: a “drinkable IV drip” to replace water and the salts the body loses"],
        items_dark=True)

content('light', '03', 'The Science',
        "Designed to taste like what your body loses.",
        [("Researchers studied ", False), ("the composition of sweat", True),
         (" and built an ", False), ("isotonic", True),
         (" drink that mirrors the body’s own fluids, so water and electrolytes absorb quickly.", False)],
        items=["Sodium, potassium, calcium & magnesium — the ions lost in perspiration",
               "Light sugars and citrus for fast absorption and a clean, mild taste",
               "Non-carbonated — made for replenishment, not just refreshment"])

content('dark', '04', 'The Long Road',
        "1,000 prototypes before it hit the shelf.",
        [("Early versions tasted ", False), ("far too bitter", True),
         (". The breakthrough came when the team added a dash of ", False),
         ("citrus juice powder", True),
         (", finally landing on two recipes with slightly different sugar levels.", False)])

content('light', '05', 'The Name',
        "“Sweat” — a marketing gamble that became the brand.",
        [("“Pocari”", True), (" means nothing — invented purely for its bright, refreshing sound. ", False),
         ("“Sweat”", True), (" bluntly names what the drink replaces.", False)],
        items=["A risky name in English-speaking markets — but unforgettable",
               "The blue-and-white can evokes a clear sky and clean water"])

content('night', '06', 'The Launch',
        "1980: a slow start, then a flood.",
        [("Sales were sluggish at first — the concept was unfamiliar. Otsuka responded by ", False),
         ("giving away millions of free samples", True),
         (", letting the taste sell itself. It worked.", False)],
        items_dark=True)

content('light', '07', 'Going Global',
        "From Tokushima to the world.",
        [("It became ", False), ("Japan’s first home-grown non-alcoholic drink to ship over $1 billion", True),
         (". Sold as cans, then bottles, plus powder and the lighter Ion Water.", False)],
        stats=[("1982", "Hong Kong & Taiwan"), ("1983", "Singapore & Mid-East"),
               ("$1B+", "by the mid-1990s"), ("20+", "countries today")])

content('dark', '08', 'Culture',
        "A drink that became a way of life in Asia.",
        [("Beyond sport, Pocari wove itself into daily life — ", False),
         ("dreamy anime and cinematic ad campaigns", True),
         (", a youthful blue identity, and staple status across ", False),
         ("Indonesia, Korea and Southeast Asia", True), (".", False)],
        items=["One of the most recognizable beverage brands in Asia",
               "Marketed for fever, exercise, travel and everyday hydration"],
        items_dark=True)

content('night', '09', 'To the Moon',
        "The first sports drink aimed at the Moon.",
        [("In 2014 Otsuka unveiled the ", False), ("Lunar Dream Capsule", True),
         (" — a Pocari-shaped titanium time capsule holding ", False),
         ("powdered Pocari Sweat", True),
         (" and children’s dreams, meant to ride a private lander to the Moon.", False)],
        items=["The dream: open it in ~30 years and mix the powder with water mined on the Moon"],
        items_dark=True)

# closing
s = prs.slides.add_slide(BLANK); bg(s, 'dark'); logo(s, True)
tf = tb(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.5))
eyebrow(tf, "The Legacy", ICE)
p = tf.add_paragraph(); p.space_before = Pt(14)
setrun(p.add_run(), "“Closer to your body than water itself.”", 44, WHITE, bold=True, italic=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(20)
setrun(p2.add_run(), "From a hospital bed in Mexico to a moonbound capsule — Pocari Sweat turned "
       "pharmaceutical science into an everyday icon, and rewrote what a drink could be.", 21, ICE, bold=False)
p3 = tf.add_paragraph(); p3.space_before = Pt(24)
setrun(p3.add_run(), "Sources: Otsuka Pharmaceutical · CNN Business · Wikipedia · Atlas Obscura", 13, ICE, bold=True)

prs.save("Pocari_Sweat_History.pptx")
print("Saved Pocari_Sweat_History.pptx with", len(prs.slides._sldIdLst), "slides")
